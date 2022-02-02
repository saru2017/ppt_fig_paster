#!/usr/bin/env python
# -*- coding: utf-8 -*-

import pptx
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from pptx.dml.color import RGBColor
import glob


def place_fig(filename, num):
    global g_slide, g_gap

    number_in_slide = num % 9
    row_in_slide = num % 3
    column_in_slide = num // 3
    position_x = 1.0 + g_gap + (6.3 + g_gap) * row_in_slide
    position_y = 1.4 + g_gap + (8.8 + g_gap) * column_in_slide
    g_slide.shapes.add_picture(filename, Cm(position_x), Cm(position_y), Cm(6.3), Cm(8.8))



ppt = pptx.Presentation()
ppt.slide_width = Cm(21.0)
ppt.slide_height = Cm(29.7)
g_gap = 0.03
card_wdith = 6.3
card_height = 8.8

filenames = glob.glob("./data/*.jpg")
print(filenames)

count = len(filenames) * 4

for i in range(count):
    filename = filenames[i // 4]
    if i % 9 == 0:
        blank_slide_layout = ppt.slide_layouts[6]
        #スライドを追加#
        g_slide = ppt.slides.add_slide(blank_slide_layout)

        shape = g_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.0), Cm(1.4), Cm(6.3 * 3 + g_gap * 3), Cm(8.8 * 3 + g_gap * 3))
        shape.line.color.rgb = RGBColor(0, 0, 203)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 0, 203)

    
    place_fig(filename, i % 9)

ppt.save("./output.pptx")

